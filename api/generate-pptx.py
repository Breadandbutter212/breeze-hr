from http.server import BaseHTTPRequestHandler
import json
import base64
import io
import copy
import os
import re
import urllib.request


def _verify_supabase_user(headers):
    """Validate a Supabase JWT via the auth REST endpoint. Returns True if the token is a valid user."""
    auth = headers.get('Authorization') or headers.get('authorization') or ''
    token = auth[7:].strip() if auth[:7].lower() == 'bearer ' else ''
    if not token:
        return False
    sb_url = (os.environ.get('SUPABASE_URL') or '').rstrip('/')
    key = os.environ.get('SUPABASE_SERVICE_KEY') or os.environ.get('SUPABASE_ANON_KEY') or ''
    if not sb_url or not key:
        return False
    try:
        req = urllib.request.Request(sb_url + '/auth/v1/user',
                                     headers={'Authorization': 'Bearer ' + token, 'apikey': key})
        with urllib.request.urlopen(req, timeout=8) as resp:
            return resp.status == 200
    except Exception:
        return False

R_EMBED = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
R_IMAGE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image'


# ── XLSX generation (openpyxl) - folded in here to stay under the serverless function limit ──
def _xl_is_formula(v):
    return isinstance(v, str) and v.strip().startswith('=')


def _xl_clean_num(s):
    return re.sub(r'[£$€,\s]', '', s)


def _xl_col_type(rows, ci):
    cur = pct = any_ = False
    num = True
    for ri in range(1, len(rows)):
        raw = '' if ci >= len(rows[ri]) or rows[ri][ci] is None else str(rows[ri][ci]).strip()
        if not raw or _xl_is_formula(raw):
            continue
        any_ = True
        if re.search(r'[£$€]', raw):
            cur = True
        if raw.endswith('%'):
            pct = True
        if not re.match(r'^-?[\d,]+(\.\d+)?$', re.sub(r'[£$€%\s]', '', raw)):
            num = False
    if not any_:
        return 'text'
    if pct and num:
        return 'percent'
    if cur and num:
        return 'currency'
    if num:
        return 'number'
    return 'text'


# House chart palette (navy-led, tasteful)
XL_PALETTE = ['#2A5080', '#10B981', '#F59E0B', '#8B5CF6', '#EF4444', '#0EA5E9', '#14B8A6', '#6366F1']
XL_FONT = 'Calibri'


def _xl_axis_numfmt(types, cols):
    ct = set(types[c] for c in cols)
    if ct == {'percent'}:
        return '0%'
    if ct == {'currency'}:
        return u'£#,##0'
    return '#,##0'


def _xl_add_chart(wb, ws, name, rows, types, max_cols, chart_type, base_row=0, base_col=0, anchor=None):
    """One house-styled chart for a table whose top-left header cell is at (base_row, base_col)."""
    num_cols = [c for c in range(1, max_cols) if types[c] in ('number', 'currency', 'percent')]
    if not num_cols:
        return
    n = len(rows)
    last = n
    last_label = str(rows[-1][0]).strip().lower() if rows and rows[-1] and rows[-1][0] is not None else ''
    if re.match(r'^(total|subtotal|grand total|average|mean)\b', last_label) and n - 1 >= 3:
        last = n - 1
    fd, ld = base_row + 1, base_row + last - 1  # absolute data-row range
    bc = base_col

    hdr = [('' if c is None else str(c)).strip() for c in rows[0]]
    time_rx = re.compile(r'^(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec|q[1-4]|fy|h[12]|20\d\d)', re.I)
    tcount = sum(1 for h in hdr[1:] if time_rx.match(h))
    is_time_matrix = (max_cols - 1) >= 3 and tcount >= (max_cols - 1) * 0.6 and len(num_cols) >= 3

    grid = {'visible': True, 'line': {'color': '#EDF1F6'}}
    num_font = {'name': XL_FONT, 'size': 9, 'color': '#64748B'}
    legend_font = {'name': XL_FONT, 'size': 9}
    yfmt = _xl_axis_numfmt(types, num_cols)

    if is_time_matrix:
        chart = wb.add_chart({'type': 'line'})
        for i, ri in enumerate(range(base_row + 1, base_row + last)):
            color = XL_PALETTE[i % len(XL_PALETTE)]
            chart.add_series({
                'name': [name, ri, bc],
                'categories': [name, base_row, bc + 1, base_row, bc + max_cols - 1],
                'values': [name, ri, bc + 1, ri, bc + max_cols - 1],
                'line': {'color': color, 'width': 2.0},
                'marker': {'type': 'circle', 'size': 4, 'border': {'color': color}, 'fill': {'color': color}},
            })
        chart.set_x_axis({'num_font': num_font, 'line': {'color': '#CBD5E1'}})
        chart.set_y_axis({'num_format': yfmt, 'num_font': num_font, 'major_gridlines': grid, 'line': {'none': True}})
        chart.set_legend({'position': 'bottom', 'font': legend_font})
    elif chart_type == 'pie':
        chart = wb.add_chart({'type': 'pie'})
        pts = [{'fill': {'color': XL_PALETTE[i % len(XL_PALETTE)]}} for i in range(ld - fd + 1)]
        chart.add_series({
            'name': [name, base_row, bc + num_cols[0]],
            'categories': [name, fd, bc, ld, bc],
            'values': [name, fd, bc + num_cols[0], ld, bc + num_cols[0]],
            'points': pts,
            'data_labels': {'percentage': True, 'font': {'name': XL_FONT, 'size': 9, 'color': '#FFFFFF', 'bold': True}},
        })
        chart.set_legend({'position': 'right', 'font': legend_font})
    else:
        ctype = 'bar' if chart_type == 'bar' else ('line' if chart_type == 'line' else 'column')
        chart = wb.add_chart({'type': ctype})
        single = len(num_cols) == 1
        for i, c in enumerate(num_cols):
            color = XL_PALETTE[i % len(XL_PALETTE)]
            ser = {
                'name': [name, base_row, bc + c],
                'categories': [name, fd, bc, ld, bc],
                'values': [name, fd, bc + c, ld, bc + c],
            }
            if ctype == 'line':
                ser['line'] = {'color': color, 'width': 2.0}
                ser['marker'] = {'type': 'circle', 'size': 4, 'border': {'color': color}, 'fill': {'color': color}}
            else:
                ser['fill'] = {'color': color}
                ser['border'] = {'none': True}
            if single:
                ser['data_labels'] = {'value': True, 'num_format': yfmt, 'font': {'name': XL_FONT, 'size': 9, 'color': '#475569'}}
            chart.add_series(ser)
        chart.set_x_axis({'num_font': num_font, 'line': {'color': '#CBD5E1'}})
        chart.set_y_axis({'num_format': yfmt, 'num_font': num_font, 'major_gridlines': grid, 'line': {'none': True}})
        chart.set_legend({'none': True} if single else {'position': 'bottom', 'font': legend_font})

    chart.set_title({'name': name, 'name_font': {'name': XL_FONT, 'size': 12, 'bold': True, 'color': '#1A2433'}})
    chart.set_chartarea({'border': {'none': True}})
    chart.set_plotarea({'border': {'none': True}})
    chart.set_size({'width': 470, 'height': 280})

    if anchor is None:
        anchor = (base_row + 1, base_col + max_cols + 1) if max_cols <= 6 else (base_row + len(rows) + 2, base_col)
    ws.insert_chart(anchor[0], anchor[1], chart, {'x_offset': 8, 'y_offset': 2})


def build_workbook(title, sheets):
    import xlsxwriter

    output = io.BytesIO()
    wb = xlsxwriter.Workbook(output, {'in_memory': True})

    header_fmt = wb.add_format({'bold': True, 'font_color': '#FFFFFF', 'bg_color': '#2A5080', 'font_name': XL_FONT,
                                'font_size': 11, 'align': 'left', 'valign': 'vcenter', 'border': 1, 'border_color': '#1F3A5F'})
    base = {'font_name': XL_FONT, 'font_size': 11, 'border': 1, 'border_color': '#E2E8F0'}
    text_fmt = wb.add_format(base)
    num_fmt = wb.add_format(dict(base, num_format='#,##0'))
    cur_fmt = wb.add_format(dict(base, num_format=u'£#,##0'))
    pct_fmt = wb.add_format(dict(base, num_format='0.0%'))
    fmt_for = {'currency': cur_fmt, 'percent': pct_fmt, 'number': num_fmt}

    used_names = set()

    def safe_name(name, idx):
        n = re.sub(r'[\\/?*\[\]:]', ' ', str(name or ('Sheet%d' % (idx + 1)))).strip()[:31] or ('Sheet%d' % (idx + 1))
        base_n, k = n[:28], 2
        while n.lower() in used_names:
            n = ('%s %d' % (base_n, k))[:31]
            k += 1
        used_names.add(n.lower())
        return n

    def norm(rs):
        rs = [r if isinstance(r, list) else [r] for r in (rs or [])]
        return rs if rs else [['']]

    def write_table(ws, rows, types, br, bc):
        for ri, r in enumerate(rows):
            for c in range(len(types)):
                raw = '' if c >= len(r) or r[c] is None else str(r[c])
                R, C = br + ri, bc + c
                if ri == 0:
                    ws.write(R, C, raw, header_fmt)
                    continue
                t = types[c]
                f = fmt_for.get(t, text_fmt)
                if _xl_is_formula(raw):
                    ws.write_formula(R, C, raw, f)
                elif raw == '':
                    ws.write_blank(R, C, None, text_fmt)
                elif t == 'currency':
                    try:
                        ws.write_number(R, C, float(_xl_clean_num(raw)), cur_fmt)
                    except ValueError:
                        ws.write(R, C, raw, text_fmt)
                elif t == 'percent':
                    try:
                        ws.write_number(R, C, float(re.sub(r'[%\s,]', '', raw)) / 100.0, pct_fmt)
                    except ValueError:
                        ws.write(R, C, raw, text_fmt)
                elif t == 'number':
                    try:
                        v = _xl_clean_num(raw)
                        ws.write_number(R, C, int(v) if re.match(r'^-?\d+$', v) else float(v), num_fmt)
                    except ValueError:
                        ws.write(R, C, raw, text_fmt)
                else:
                    ws.write(R, C, raw, text_fmt)

    if not sheets:
        sheets = [{'name': title or 'Sheet1', 'rows': [['(empty)']], 'chart': None}]
    if len(sheets) == 1 and not sheets[0].get('blocks') and (not sheets[0].get('name') or sheets[0].get('name') == 'Sheet1'):
        sheets[0]['name'] = title or 'Sheet1'

    for idx, s in enumerate(sheets):
        name = safe_name(s.get('name'), idx)
        ws = wb.add_worksheet(name)
        blocks = s.get('blocks')

        if blocks:
            # Dashboard: tables stacked down the left, charts in a 2-col grid to the right
            placed = []
            cur_row = 0
            col_w = {}
            for b in blocks:
                rows = norm(b.get('rows'))
                mc = max(len(r) for r in rows)
                types = [_xl_col_type(rows, c) for c in range(mc)]
                write_table(ws, rows, types, cur_row, 0)
                placed.append({'rows': rows, 'types': types, 'mc': mc, 'br': cur_row, 'chart': (b.get('chart') or '').lower()})
                for c in range(mc):
                    longest = col_w.get(c, 9)
                    for r in rows:
                        v = '' if c >= len(r) or r[c] is None else str(r[c])
                        longest = max(longest, len(v) + 2)
                    col_w[c] = longest
                cur_row += len(rows) + 2
            table_w = max((p['mc'] for p in placed), default=1)
            for c, w in col_w.items():
                ws.set_column(c, c, min(max(w, 9), 60))
            chart_col = table_w + 1
            ci = 0
            for p in placed:
                if p['chart'] and len(p['rows']) >= 3:
                    anchor = ((ci // 2) * 16, chart_col + (ci % 2) * 9)
                    _xl_add_chart(wb, ws, name, p['rows'], p['types'], p['mc'], p['chart'], base_row=p['br'], base_col=0, anchor=anchor)
                    ci += 1
        else:
            rows = norm(s.get('rows'))
            max_cols = max(len(r) for r in rows)
            types = [_xl_col_type(rows, c) for c in range(max_cols)]
            for c in range(max_cols):
                longest = 9
                for r in rows:
                    v = '' if c >= len(r) or r[c] is None else str(r[c])
                    longest = max(longest, len(v) + 2)
                ws.set_column(c, c, min(max(longest, 9), 60))
            write_table(ws, rows, types, 0, 0)
            ws.freeze_panes(1, 0)
            if len(rows) > 1:
                ws.autofilter(0, 0, len(rows) - 1, max_cols - 1)
            chart_type = (s.get('chart') or '').lower()
            if chart_type and len(rows) >= 3:
                _xl_add_chart(wb, ws, name, rows, types, max_cols, chart_type)

    wb.close()
    return output.getvalue()

def clone_slide_with_images(prs, source_idx):
    """Clone a slide preserving images, background shapes and decorative elements."""
    from pptx.oxml.ns import qn
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)

    # Clear shapes added by add_slide
    dest_tree = new_slide.shapes._spTree
    for el in list(dest_tree):
        dest_tree.remove(el)

    rId_map = {}
    for el in source.shapes._spTree:
        tag = el.tag.split('}')[-1]
        if tag == 'pic':
            blip = el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blip is not None:
                old_rId = blip.get(R_EMBED)
                if old_rId and old_rId not in rId_map:
                    try:
                        image_part = source.part.related_part(old_rId)
                        new_rId = new_slide.part.relate_to(image_part, R_IMAGE)
                        rId_map[old_rId] = new_rId
                    except Exception:
                        pass
            new_el = copy.deepcopy(el)
            blip_new = new_el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blip_new is not None:
                old_rId = blip_new.get(R_EMBED)
                if old_rId in rId_map:
                    blip_new.set(R_EMBED, rId_map[old_rId])
            dest_tree.append(new_el)
        else:
            dest_tree.append(copy.deepcopy(el))

    return new_slide

def overwrite_title_text(ph, new_title):
    """Slide is already a clone - just swap the text in the existing runs. Touch nothing else."""
    from pptx.oxml.ns import qn
    txBody = ph.text_frame._txBody
    # Find all a:t elements in the title and set the first one, clear the rest
    all_t = txBody.findall('.//' + qn('a:t'))
    if all_t:
        all_t[0].text = new_title
        for t in all_t[1:]:
            t.text = ''
    else:
        # Fallback: no runs found - use tf.text
        ph.text_frame.text = new_title

def overwrite_body_text(ph, bullets):
    """Slide is already a clone - replace body paragraph texts. Clone first run for formatting."""
    from pptx.oxml.ns import qn
    from lxml import etree
    if not bullets:
        return
    txBody = ph.text_frame._txBody
    paras = txBody.findall(qn('a:p'))
    if not paras:
        return

    # Use first existing paragraph as formatting template
    ref_para = copy.deepcopy(paras[0])

    # Remove all existing paragraphs
    for p in paras:
        txBody.remove(p)

    for bullet in bullets:
        new_p = copy.deepcopy(ref_para)
        # Set text in first run, clear rest
        all_t = new_p.findall('.//' + qn('a:t'))
        if all_t:
            all_t[0].text = bullet
            for t in all_t[1:]:
                t.text = ''
        txBody.append(new_p)

def apply_ops(prs, ops):
    from pptx.oxml.ns import qn
    from pptx.enum.shapes import PP_PLACEHOLDER
    from lxml import etree

    TITLE_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)

    for op in ops:
        op_type = op.get('op', '')
        # 'slide' is a number for add_bullet/update_* but a dict for insert_slide - guard carefully
        raw_slide = op.get('slide', 1)
        slide_num = raw_slide if isinstance(raw_slide, int) else 1
        slide_idx = slide_num - 1

        # ── ADD BULLET ──────────────────────────────────────────────────────────
        if op_type == 'add_bullet' and 0 <= slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            bullet_text = op.get('bullet', '')
            after_text  = op.get('after_text', '').strip().lower()

            best_tf      = None
            best_para_idx = None

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame

                if after_text:
                    for i, para in enumerate(tf.paragraphs):
                        if after_text in para.text.strip().lower():
                            best_tf = tf
                            best_para_idx = i
                            break
                    if best_tf:
                        break
                else:
                    # Pick the frame with the most content paragraphs
                    content_count = sum(1 for p in tf.paragraphs if p.text.strip())
                    if best_tf is None or content_count > sum(1 for p in best_tf.paragraphs if p.text.strip()):
                        best_tf = tf

            if best_tf is None:
                continue

            # Clone the reference paragraph for formatting
            paras = best_tf.paragraphs
            ref_idx = best_para_idx if best_para_idx is not None else len(paras) - 1
            ref_p = paras[ref_idx]._p

            new_p = copy.deepcopy(ref_p)

            # Clear all runs from the clone, keep pPr
            for r in new_p.findall(qn('a:r')):
                new_p.remove(r)
            for br in new_p.findall(qn('a:br')):
                new_p.remove(br)

            # Build a run using the first run's rPr from the reference (preserves font/size)
            ref_runs = ref_p.findall(qn('a:r'))
            if ref_runs:
                new_r = copy.deepcopy(ref_runs[0])
                t = new_r.find(qn('a:t'))
                if t is None:
                    t = etree.SubElement(new_r, qn('a:t'))
                t.text = bullet_text
            else:
                new_r = etree.SubElement(new_p, qn('a:r'))
                t = etree.SubElement(new_r, qn('a:t'))
                t.text = bullet_text

            new_p.append(new_r)

            # Insert after the reference paragraph
            ref_p.addnext(new_p)

        # ── INSERT SLIDE ─────────────────────────────────────────────────────────
        elif op_type == 'insert_slide':
            after_num  = op.get('after', len(prs.slides))
            slide_data = op.get('slide', {})
            title_text = slide_data.get('title', 'New Slide')
            bullets    = slide_data.get('bullets', [])

            # Use the deck's slide layouts (inherit theme/fonts) with proper placeholders
            # This works for ALL decks including those with no placeholder-based slides
            target_layout = prs.slide_layouts[1]  # Title and Content layout
            for layout in prs.slide_layouts:
                has_title = any(ph.placeholder_format.type in TITLE_TYPES for ph in layout.placeholders)
                has_body  = any(ph.placeholder_format.idx == 1 for ph in layout.placeholders)
                if has_title and has_body:
                    target_layout = layout
                    break

            new_slide = prs.slides.add_slide(target_layout)
            for ph in new_slide.placeholders:
                ph_type = ph.placeholder_format.type
                ph_idx  = ph.placeholder_format.idx
                if ph_type in TITLE_TYPES:
                    ph.text = title_text
                elif ph_idx == 1 and bullets:
                    tf = ph.text_frame
                    tf.clear()
                    for i, b in enumerate(bullets):
                        if i == 0:
                            tf.paragraphs[0].text = b
                        else:
                            p = tf.add_paragraph()
                            p.text = b

            # Move to correct position (add_slide always appends)
            xml_slides = prs.slides._sldIdLst
            slides_list = list(xml_slides)
            new_el = slides_list[-1]
            xml_slides.remove(new_el)
            insert_pos = min(int(after_num), len(list(xml_slides)))
            xml_slides.insert(insert_pos, new_el)

        # ── UPDATE TITLE ──────────────────────────────────────────────────────────
        elif op_type == 'update_title' and 0 <= slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            new_title = op.get('title', '')
            for ph in slide.placeholders:
                if ph.placeholder_format.type in TITLE_TYPES:
                    ph.text = new_title
                    break

        # ── UPDATE BULLETS ────────────────────────────────────────────────────────
        elif op_type == 'update_bullets' and 0 <= slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            new_bullets = op.get('bullets', [])
            # Find body placeholder
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    tf = ph.text_frame
                    tf.clear()
                    for i, b in enumerate(new_bullets):
                        if i == 0:
                            tf.paragraphs[0].text = b
                        else:
                            p = tf.add_paragraph()
                            p.text = b
                    break

        # ── UPDATE DATES ──────────────────────────────────────────────────────────
        elif op_type == 'update_dates':
            import re
            year = str(op.get('year', ''))
            if year:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    run.text = re.sub(r'\b20\d{2}\b', year, run.text)


# ── GENERATION HELPERS (Claude.ai verified approach) ────────────────────────
import math

class BRAND:
    DARK   = None; DARKER = None; ACCENT = None
    TEXT   = None; MUTED  = None; FAINT  = None
    WHITE  = None; LIGHT  = None; LIGHTTXT = None
    HEAD_FONT = "Trebuchet MS"; BODY_FONT = "Calibri"

    @classmethod
    def init(cls, accent_hex, dark_hex):
        from pptx.dml.color import RGBColor
        def rgb(h): h=h.lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
        cls.ACCENT   = rgb(accent_hex or '0D9488')
        cls.DARK     = rgb(dark_hex   or '1E293B')
        cls.DARKER   = rgb('0F172A')
        cls.TEXT     = cls.DARK
        cls.MUTED    = rgb('64748B')
        cls.FAINT    = rgb('94A3B8')
        cls.LIGHT    = rgb('F8FAFC')
        cls.WHITE    = rgb('FFFFFF')
        cls.LIGHTTXT = rgb('CBD5E1')

SLIDE_W = 13.333; SLIDE_H = 7.5; MARGIN = 0.7
CONTENT_W = SLIDE_W - 2 * MARGIN

def _blank_slide(prs, bg):
    from pptx.enum.shapes import MSO_SHAPE
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fill = sl.background.fill; fill.solid(); fill.fore_color.rgb = bg
    return sl

def _add_rect(sl, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s

def _estimate_lines(text, box_w, font_pt, avg=0.50):
    cw = (font_pt * avg) / 72.0
    cpl = max(1, int(box_w / cw))
    return sum(math.ceil(max(1,len(l))/cpl) for l in text.split('\n'))

def _fit_font(text, bw, bh, start, min_pt=10, ls=1.18):
    pt = start
    while pt > min_pt:
        if _estimate_lines(text, bw, pt) * (pt*ls/72.0) <= bh: break
        pt -= 1
    return pt

def _add_text(sl, text, x, y, w, h, *, size=18, color, font=None, bold=False,
              align=None, anchor=None, autofit=True):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    font = font or BRAND.BODY_FONT
    bx = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    if autofit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        size = _fit_font(text, w, h, size)
    if anchor: tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, Pt(0))
    p = tf.paragraphs[0]
    if align: p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.name = font; run.font.color.rgb = color
    return bx

def _set_bullet_xml(paragraph):
    from pptx.oxml.ns import qn
    from pptx.util import Pt
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('marL', str(int(Pt(0.28*72).emu // 1)))
    pPr.set('indent', str(-int(Pt(0.28*72).emu // 1)))
    bf = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    bc = pPr.makeelement(qn('a:buChar'), {'char': '•'})
    pPr.append(bf); pPr.append(bc)

def _add_bullets(sl, items, x, y, w, h, *, size=15, color, gap=10):
    from pptx.util import Inches, Pt
    bx = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, Pt(0))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        sz = _fit_font(item, w-0.3, h/max(1,len(items)), size)
        run = p.add_run(); run.text = item
        run.font.size = Pt(sz); run.font.name = BRAND.BODY_FONT; run.font.color.rgb = color
        _set_bullet_xml(p)
    return bx

def _title_block(sl, title, kicker=None):
    sq_y = 0.82 if kicker else 0.66
    _add_rect(sl, MARGIN, sq_y, 0.18, 0.18, BRAND.ACCENT)
    if kicker:
        _add_text(sl, kicker.upper(), MARGIN+0.3, 0.5, 11, 0.3,
                  size=11, color=BRAND.ACCENT, bold=True, autofit=False)
    _add_text(sl, title, MARGIN+0.28, (0.74 if kicker else 0.5), CONTENT_W, 0.7,
              size=28, color=BRAND.DARK, font=BRAND.HEAD_FONT, bold=True)

def generate_deck(slides, accent_hex, dark_hex, prs_title):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    BRAND.init(accent_hex, dark_hex)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W); prs.slide_height = Inches(SLIDE_H)

    for s in slides:
        st      = s.get('type','content')
        title   = s.get('title','')
        bullets = s.get('bullets',[])
        subtitle = s.get('subtitle','')

        if st == 'title':
            sl = _blank_slide(prs, BRAND.DARKER)
            _add_rect(sl, 0.9, 1.05, 0.26, 0.26, BRAND.ACCENT)
            _add_text(sl, title, 0.88, 2.25, 11.4, 1.5, size=50,
                      color=BRAND.WHITE, font=BRAND.HEAD_FONT, bold=True)
            if subtitle:
                _add_text(sl, subtitle, 0.9, 3.95, 10.5, 0.7, size=18, color=BRAND.LIGHTTXT)

        elif st == 'section':
            sl = _blank_slide(prs, BRAND.DARK)
            _add_text(sl, '—', 0.85, 2.0, 4.0, 2.2, size=120,
                      color=BRAND.ACCENT, font=BRAND.HEAD_FONT, bold=True, autofit=False)
            _add_rect(sl, 5.1, 2.55, 0.05, 1.7, BRAND.ACCENT)
            _add_text(sl, title, 5.5, 2.7, 7.0, 0.8, size=34,
                      color=BRAND.WHITE, font=BRAND.HEAD_FONT, bold=True)

        elif st == 'close':
            sl = _blank_slide(prs, BRAND.DARKER)
            _add_rect(sl, 0.9, 1.15, 0.26, 0.26, BRAND.ACCENT)
            _add_text(sl, title or 'Questions?', 0.88, 1.6, 11, 1.0, size=44,
                      color=BRAND.WHITE, font=BRAND.HEAD_FONT, bold=True)
            if bullets:
                _add_bullets(sl, bullets, MARGIN+0.3, 3.0, 9.5, 2.2, size=17, color=BRAND.LIGHTTXT)

        elif st == 'image':
            sl = _blank_slide(prs, BRAND.WHITE)
            _title_block(sl, title, kicker='SECTION')
            _add_rect(sl, MARGIN+0.25, 1.85, 5.8, 4.4, RGBColor(0xF1,0xF5,0xF9))
            hint = s.get('imageHint','Add image here')
            _add_text(sl, f'\U0001f4f7  {hint}', MARGIN+0.25, 3.85, 5.8, 0.4,
                      size=12, color=BRAND.MUTED, align=PP_ALIGN.CENTER, autofit=False)
            if bullets:
                _add_bullets(sl, bullets, 7.3, 1.85, 5.3, 4.4, size=15, color=BRAND.TEXT)

        elif st == 'chart':
            sl = _blank_slide(prs, BRAND.WHITE)
            _title_block(sl, title)
            labels = s.get('labels',[]); series = s.get('series',[])
            ct_map = {'bar':XL_CHART_TYPE.COLUMN_CLUSTERED,'barh':XL_CHART_TYPE.BAR_CLUSTERED,
                      'line':XL_CHART_TYPE.LINE,'pie':XL_CHART_TYPE.PIE,'doughnut':XL_CHART_TYPE.DOUGHNUT}
            ct = ct_map.get(s.get('chartType','bar'), XL_CHART_TYPE.COLUMN_CLUSTERED)
            if labels and series:
                cd = ChartData(); cd.categories = labels
                for ser in series:
                    vals = [float(v) if str(v).replace('.','').replace('-','').isdigit() else 0 for v in ser.get('values',[])]
                    cd.add_series(ser.get('name',''), vals)
                chart = sl.shapes.add_chart(ct, Inches(MARGIN), Inches(1.6), Inches(CONTENT_W), Inches(5.0), cd).chart
                chart.has_legend = len(series) > 1
            note = s.get('note','')
            if note: _add_text(sl, note, MARGIN, 6.9, CONTENT_W, 0.3, size=11, color=BRAND.MUTED, autofit=False)

        else:  # content
            sl = _blank_slide(prs, BRAND.WHITE)
            _title_block(sl, title, kicker='SECTION')
            if bullets:
                _add_bullets(sl, bullets, MARGIN+0.25, 1.85, 6.0, 4.6, size=15, color=BRAND.TEXT)

    return prs

def _replace_text_in_shape(shape, new_text):
    if not shape.has_text_frame: return
    from pptx.oxml.ns import qn
    txBody = shape.text_frame._txBody
    all_t = txBody.findall('.//' + qn('a:t'))
    if all_t:
        all_t[0].text = new_text
        for t in all_t[1:]: t.text = ''

def generate_deck_from_template(slides, template_b64):
    """Generate a new deck by cloning slides from a brand template.
    Template slide mapping (by index):
      0 = title/close  (dark background)
      1 = content      (first content slide)
      2 = section      (section divider)
      3 = image        (bullets + image)
      4 = chart
    Falls back to index 1 for unrecognised types.
    """
    from pptx import Presentation
    template_bytes = base64.b64decode(template_b64)
    template_prs = Presentation(io.BytesIO(template_bytes))
    n_tmpl = len(template_prs.slides)

    TYPE_MAP = {
        'title':   min(0, n_tmpl - 1),
        'close':   min(0, n_tmpl - 1),
        'section': min(2, n_tmpl - 1),
        'image':   min(3, n_tmpl - 1),
        'chart':   min(4, n_tmpl - 1),
        'content': min(3, n_tmpl - 1),  # use Bullets slide, not Agenda slide
    }
    DEFAULT_IDX = min(3, n_tmpl - 1)

    prs = Presentation(io.BytesIO(template_bytes))
    original_count = len(prs.slides)

    for slide_data in slides:
        stype   = slide_data.get('type', 'content')
        title   = slide_data.get('title', '')
        bullets = slide_data.get('bullets', [])
        tmpl_idx = TYPE_MAP.get(stype, DEFAULT_IDX)

        new_slide = clone_slide_with_images(prs, tmpl_idx)
        _inject_template_slide(new_slide, title, bullets)

    # Remove original template slides
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    for i in range(original_count):
        s = slides_list[i]
        r_id = s.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        try:
            prs.part.drop_rel(r_id)
        except Exception:
            pass
        xml_slides.remove(s)

    return prs

PPTX_MIME = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
SKILL_BETAS = ['skills-2025-10-02', 'code-execution-2025-08-25', 'files-api-2025-04-14']
SKILL_TOOLS = [{'type': 'code_execution_20250825', 'name': 'code_execution'}]
SKILL_DEF   = [{'type': 'anthropic', 'skill_id': 'pptx', 'version': 'latest'}]

def _upload_template(client, template_b64):
    """Upload template bytes to Files API, return file_id."""
    import tempfile, os
    tmpl_bytes = base64.b64decode(template_b64)
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        f.write(tmpl_bytes); tmp = f.name
    try:
        with open(tmp, 'rb') as f:
            obj = client.beta.files.upload(file=('template.pptx', f, PPTX_MIME))
        return obj.id
    finally:
        os.unlink(tmp)

def _extract_file_id(message):
    """Find file_id of generated pptx in tool result blocks."""
    file_id = None
    for block in message.content:
        t = getattr(block, 'type', '')
        if t in ('code_execution_tool_result', 'bash_code_execution_tool_result'):
            result = getattr(block, 'content', None)
            if getattr(result, 'type', '') in ('code_execution_result', 'bash_code_execution_result'):
                for output in getattr(result, 'content', []) or []:
                    if getattr(output, 'file_id', None):
                        file_id = output.file_id
    return file_id

def generate_with_skill(topic, instructions=None, template_b64=None):
    """Call Anthropic pptx Agent Skill to generate a deck. Returns bytes."""
    import anthropic
    client = anthropic.Anthropic(api_key=os.environ.get('ANTHROPIC_API_KEY'))

    # Use stored template file_id or upload fresh
    template_file_id = os.environ.get('TEMPLATE_FILE_ID')
    if not template_file_id and template_b64:
        template_file_id = _upload_template(client, template_b64)

    prompt = (
        f'Create a brand-new PowerPoint presentation from scratch about:\n{topic}\n\n'
        'There is NO input file and NO template. Do not search for, open, or try '
        'to load any existing .pptx, and do NOT use the template-based editing '
        'workflow — build the slides yourself using the skill\'s from-scratch '
        'creation path. Pick a clean, professional design. Keep each slide '
        'focused, ensure no text overflows its boxes, and save the result as a '
        '.pptx file.'
    )
    if instructions:
        prompt += f'\n\nAdditional instructions:\n{instructions}'

    # Plain content - no cache_control on user message (it changes every request, never hits cache)
    content = [{'type': 'text', 'text': prompt}]
    if template_file_id:
        content.append({'type': 'container_upload', 'file_id': template_file_id})

    # Automatic caching: cache_control at top level caches the growing conversation
    # prefix across the skill's internal multi-turn code execution loop
    betas_with_cache = SKILL_BETAS + ['prompt-caching-2024-07-31']

    # Try streaming first (newer SDK), fall back to create() for older SDK
    try:
        with client.beta.messages.stream(
            model='claude-sonnet-4-6',
            max_tokens=8000,
            cache_control={'type': 'ephemeral'},
            betas=betas_with_cache,
            container={'skills': SKILL_DEF},
            tools=SKILL_TOOLS,
            messages=[{'role': 'user', 'content': content}]
        ) as stream:
            message = stream.get_final_message()
    except AttributeError:
        message = client.beta.messages.create(
            model='claude-sonnet-4-6',
            max_tokens=8000,
            cache_control={'type': 'ephemeral'},
            betas=betas_with_cache,
            container={'skills': SKILL_DEF},
            tools=SKILL_TOOLS,
            messages=[{'role': 'user', 'content': content}]
        )

    out_id = _extract_file_id(message)
    if not out_id:
        txt = ''.join(getattr(b,'text','') for b in message.content if getattr(b,'type','')=='text')
        raise RuntimeError(f'No output file produced. Model: {txt[:400]}')

    downloaded = client.beta.files.download(file_id=out_id)
    return downloaded.read()

class handler(BaseHTTPRequestHandler):
    def log_message(self, format, *args):
        pass  # suppress request logging

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type, Authorization')
        self.end_headers()

    def do_POST(self):
        try:
            content_length = int(self.headers.get('Content-Length', 0))
            body = self.rfile.read(content_length)
            data = json.loads(body)

            # Excel export (openpyxl) - folded into this function to stay under the serverless limit
            if data.get('action') == 'xlsx':
                xlsx_bytes = build_workbook(data.get('title', 'Spreadsheet'), data.get('sheets', []))
                self.send_response(200)
                self.send_header('Content-Type', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
                self.send_header('Content-Disposition', 'attachment; filename="spreadsheet.xlsx"')
                self.send_header('Content-Length', str(len(xlsx_bytes)))
                self.send_header('Access-Control-Allow-Origin', '*')
                self.end_headers()
                self.wfile.write(xlsx_bytes)
                return

            from pptx import Presentation

            pptx_b64 = data.get('pptxBase64', '')
            ops      = data.get('ops', [])
            action   = data.get('action', 'edit' if pptx_b64 else 'generate')

            if action == 'generate':
                slides        = data.get('slides', [])
                accent        = data.get('accent', '475569')
                dark          = data.get('dark',   '1E293B')
                prs_title     = data.get('title', 'Presentation')
                template_b64  = data.get('templateBase64', '')
                if template_b64:
                    # Extract brand colours from template, use them in generate_deck
                    try:
                        from pptx import Presentation as _Prs
                        from pptx.oxml.ns import qn as _qn
                        _tmpl = _Prs(io.BytesIO(base64.b64decode(template_b64)))
                        _theme = _tmpl.slide_master.theme_color_map
                        # Try to get accent1 and dk1 from theme XML
                        _theme_xml = _tmpl.slide_master.part.part_related_by(
                            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'
                        )._element
                        def _hex(el, tag):
                            n = el.find('.//' + _qn(tag))
                            if n is not None:
                                clr = n.find(_qn('a:srgbClr'))
                                if clr is not None: return clr.get('val','')
                            return ''
                        tmpl_accent = _hex(_theme_xml, 'a:accent1') or accent
                        tmpl_dark   = _hex(_theme_xml, 'a:dk2')     or _hex(_theme_xml, 'a:dk1') or dark
                        # Also scan slide shapes for most common background colour
                        for sl in _tmpl.slides[:3]:
                            for sh in sl.shapes:
                                try:
                                    rgb = sh.fill.fore_color.rgb
                                    val = str(rgb)
                                    # If it looks like a dark background colour use it
                                    if int(val[:2],16) < 50 and int(val[2:4],16) < 50:
                                        tmpl_dark = val
                                        break
                                except: pass
                        prs = generate_deck(slides, tmpl_accent, tmpl_dark, prs_title)
                    except Exception as _e:
                        prs = generate_deck(slides, accent, dark, prs_title)
                else:
                    prs = generate_deck(slides, accent, dark, prs_title)
            elif action == 'generate_ai':
                # Spends the server's Anthropic key, so require an authenticated caller
                if not _verify_supabase_user(self.headers):
                    self._error(401, 'Unauthorized')
                    return
                # Use Anthropic pptx Agent Skill for high-quality generation
                topic        = data.get('topic', '')
                instructions = data.get('instructions')
                template_b64 = data.get('templateBase64', '')
                if not topic:
                    self._error(400, 'Missing topic')
                    return
                pptx_bytes = generate_with_skill(topic, instructions, template_b64)
                self.send_response(200)
                self.send_header('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
                self.send_header('Content-Disposition', 'attachment; filename="presentation.pptx"')
                self.send_header('Content-Length', str(len(pptx_bytes)))
                self.send_header('Access-Control-Allow-Origin', '*')
                self.end_headers()
                self.wfile.write(pptx_bytes)
                return

            else:
                if not pptx_b64:
                    self._error(400, 'Missing pptxBase64')
                    return
                pptx_bytes = base64.b64decode(pptx_b64)
                prs = Presentation(io.BytesIO(pptx_bytes))
                apply_ops(prs, ops)

            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            result_b64 = base64.b64encode(output.read()).decode()

            self._json(200, {'pptxBase64': result_b64})

        except Exception as e:
            self._error(500, str(e))

    def _json(self, code, obj):
        body = json.dumps(obj).encode()
        self.send_response(code)
        self.send_header('Content-Type', 'application/json')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Content-Length', str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _error(self, code, msg):
        self._json(code, {'error': msg})
